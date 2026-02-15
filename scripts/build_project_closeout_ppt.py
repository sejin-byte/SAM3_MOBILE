#!/usr/bin/env python3
"""Build SAM3_MOBILE project closeout presentation (.pptx).

Output:
  docs/ppt/SAM3_MOBILE_Project_Closeout_20260215_ko.pptx
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "docs" / "ppt"
OUT_PATH = OUT_DIR / "SAM3_MOBILE_Project_Closeout_20260215_ko.pptx"

TITLE_COLOR = RGBColor(22, 32, 64)
BODY_COLOR = RGBColor(40, 40, 40)
ACCENT = RGBColor(0, 104, 179)
ACCENT_SOFT = RGBColor(224, 242, 255)
SUCCESS = RGBColor(16, 124, 16)
WARNING = RGBColor(176, 108, 0)
ERROR = RGBColor(185, 28, 28)

FONT_KO = "Apple SD Gothic Neo"


def set_run_style(run, size=18, bold=False, color=BODY_COLOR):
    run.font.name = FONT_KO
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_paragraph_style(paragraph, size=18, bold=False, color=BODY_COLOR):
    if not paragraph.runs:
        paragraph.add_run()
    for run in paragraph.runs:
        set_run_style(run, size=size, bold=bold, color=color)


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    set_paragraph_style(p, size=32, bold=True, color=TITLE_COLOR)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.2), Inches(0.55))
        tf2 = sub_box.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        set_paragraph_style(p2, size=16, color=RGBColor(90, 90, 90))

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.72), Inches(12.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_bullets(slide, left, top, width, height, bullets: Iterable[str], size=20):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(8)
        set_paragraph_style(p, size=size, color=BODY_COLOR)


def add_kpi_card(slide, left, top, width, height, title, value, note, color=ACCENT_SOFT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1.5)

    tbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), width - Inches(0.3), Inches(0.3))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_style(p, size=14, bold=True, color=TITLE_COLOR)

    vbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.45), width - Inches(0.3), Inches(0.4))
    vtf = vbox.text_frame
    vtf.clear()
    vp = vtf.paragraphs[0]
    vp.text = value
    set_paragraph_style(vp, size=24, bold=True, color=ACCENT)

    nbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.9), width - Inches(0.3), Inches(0.35))
    ntf = nbox.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    np.text = note
    set_paragraph_style(np, size=12, color=RGBColor(80, 80, 80))


def add_box(slide, left, top, width, height, text, fill=RGBColor(240, 246, 255), border=ACCENT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_style(p, size=14, bold=True, color=TITLE_COLOR)
    return box


def add_arrow(slide, left, top, width, height):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arr.fill.solid()
    arr.fill.fore_color.rgb = ACCENT
    arr.line.fill.background()
    return arr


def add_timeline_block(slide, left, top, width, height, date_text, label, color):
    block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    block.fill.solid()
    block.fill.fore_color.rgb = color
    block.line.color.rgb = RGBColor(120, 120, 120)

    tf = block.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = date_text
    p1.alignment = PP_ALIGN.CENTER
    set_paragraph_style(p1, size=12, bold=True, color=TITLE_COLOR)

    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = PP_ALIGN.CENTER
    set_paragraph_style(p2, size=13, color=BODY_COLOR)


def add_bar_comparison(slide, left, top, width, height, labels, values, colors, title):
    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_style(p, size=14, bold=True, color=TITLE_COLOR)

    chart_top = top + Inches(0.4)
    chart_h = height - Inches(0.45)
    n = len(values)
    max_v = max(values) if values else 1.0
    col_w = width / max(n, 1)
    bar_w = col_w * 0.55

    axis = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, left + Inches(0.1), chart_top + chart_h, left + width - Inches(0.1), chart_top + chart_h)
    axis.line.color.rgb = RGBColor(120, 120, 120)
    axis.line.width = Pt(1)

    for idx, (label, value, color) in enumerate(zip(labels, values, colors)):
        x = left + col_w * idx + (col_w - bar_w) / 2
        h = (value / max_v) * (chart_h - Inches(0.3))
        y = chart_top + chart_h - h
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        val_box = slide.shapes.add_textbox(x - Inches(0.03), y - Inches(0.2), bar_w + Inches(0.06), Inches(0.18))
        vtf = val_box.text_frame
        vtf.clear()
        vp = vtf.paragraphs[0]
        vp.text = f"{value:.4f}"
        vp.alignment = PP_ALIGN.CENTER
        set_paragraph_style(vp, size=10, bold=True, color=BODY_COLOR)

        lb = slide.shapes.add_textbox(x - Inches(0.1), chart_top + chart_h + Inches(0.03), bar_w + Inches(0.2), Inches(0.25))
        ltf = lb.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lp.text = label
        lp.alignment = PP_ALIGN.CENTER
        set_paragraph_style(lp, size=10, color=BODY_COLOR)


def add_picture_or_placeholder(slide, path: Path, left, top, width, height, caption: str):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = RGBColor(170, 170, 170)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"이미지 없음\n{path.name}"
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_style(p, size=12, color=WARNING)

    cap = slide.shapes.add_textbox(left, top + height + Inches(0.05), width, Inches(0.22))
    ctf = cap.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.text = caption
    cp.alignment = PP_ALIGN.CENTER
    set_paragraph_style(cp, size=11, color=BODY_COLOR)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1) Cover
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(246, 250, 255)
    bg.line.fill.background()
    add_title(s, "SAM3_MOBILE 프로젝트 종료 보고", "모델 아키텍처, 원리, 실험 이력 및 배포 정리 (2026-02-15)")
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.1),
        Inches(7.4),
        Inches(2.2),
        [
            "발표 목적: 기술/비기술 이해관계자 모두가 현재 상태를 한 번에 이해",
            "범위: 배경 -> 아키텍처/원리 -> 실험 성과 -> 배포 상태 -> 다음 사이클",
            "기준 태그: v2026.02.15-project-close",
        ],
        size=20,
    )
    add_picture_or_placeholder(
        s,
        ROOT / "docs" / "visual_qa_segment_everything.png",
        Inches(8.3),
        Inches(1.9),
        Inches(4.6),
        Inches(3.8),
        "시각 결과 예시 (segment everything)",
    )

    # 2) 한 줄 요약
    s = prs.slides.add_slide(blank)
    add_title(s, "프로젝트 한 줄 요약", "SAM3 Teacher를 모바일 Student로 압축하면서 품질을 회복한 실험")
    add_kpi_card(s, Inches(0.8), Inches(2.0), Inches(3.8), Inches(1.4), "mIoU", "0.0570 -> 0.0898", "+57.46%")
    add_kpi_card(s, Inches(4.9), Inches(2.0), Inches(3.8), Inches(1.4), "Teacher 대비", "53.9%", "student/teacher mIoU", color=RGBColor(237, 250, 242))
    add_kpi_card(s, Inches(9.0), Inches(2.0), Inches(3.4), Inches(1.4), "최종 상태", "학습 종료", "배포/아카이브 단계", color=RGBColor(255, 247, 230))
    add_bullets(
        s,
        Inches(0.9),
        Inches(4.1),
        Inches(11.8),
        Inches(2.6),
        [
            "구조 안정화와 재학습을 통해 품질 하락 구간에서 회복세를 확보",
            "다만 공개 출시 품질까지는 추가 사이클(teacher 추종/온디바이스 KPI)이 필요",
            "이번 사이클은 실험 종료, 코드/문서/릴리즈 태그까지 정리 완료",
        ],
        size=20,
    )

    # 3) 배경
    s = prs.slides.add_slide(blank)
    add_title(s, "프로젝트 배경", "왜 SAM3를 모바일용으로 다시 학습했는가")
    add_bullets(
        s,
        Inches(0.8),
        Inches(1.9),
        Inches(12.0),
        Inches(4.8),
        [
            "기존 Teacher 모델은 품질은 높지만 모바일 실시간 서비스에는 계산량/크기 부담이 큼",
            "목표는 경량 Student 모델로 교체하면서 품질 손실을 최소화하는 것",
            "초기 구간에서 tokenizer/정렬 이슈로 성능 정체가 발생해 파이프라인 재정비가 필요했음",
            "결과적으로 From-scratch 재학습 + 증류 정렬로 회복 경로를 확보",
        ],
        size=21,
    )

    # 4) 목표/KPI
    s = prs.slides.add_slide(blank)
    add_title(s, "프로젝트 목표 및 KPI", "품질, 속도, 배포 가능성의 균형")
    add_box(s, Inches(0.9), Inches(2.0), Inches(3.8), Inches(1.3), "품질\n- mIoU 상승\n- teacher 격차 축소", fill=RGBColor(239, 248, 255))
    add_box(s, Inches(4.8), Inches(2.0), Inches(3.8), Inches(1.3), "효율\n- 추론 속도 유지\n- 양자화 안정성", fill=RGBColor(238, 249, 240))
    add_box(s, Inches(8.7), Inches(2.0), Inches(3.8), Inches(1.3), "배포\n- iOS/Android export\n- 운영 문서화", fill=RGBColor(255, 247, 231))
    add_bullets(
        s,
        Inches(0.9),
        Inches(4.0),
        Inches(11.8),
        Inches(2.4),
        [
            "Target KPI(운영 기준): model size < 50MB, image latency < 100ms, video fps > 15",
            "이번 사이클에서 품질 회복은 달성했지만, 온디바이스 KPI는 아직 추가 검증 필요",
        ],
        size=19,
    )

    # 5) 타임라인
    s = prs.slides.add_slide(blank)
    add_title(s, "진행 타임라인", "문제 진단 -> 구조 개선 -> 재학습 -> 종료 정리")
    add_timeline_block(s, Inches(0.8), Inches(2.3), Inches(2.3), Inches(1.2), "2026-02-12", "Baseline 측정\n(mIoU 0.057)", RGBColor(245, 245, 245))
    add_timeline_block(s, Inches(3.4), Inches(2.3), Inches(2.3), Inches(1.2), "2026-02-13", "파이프라인\n수정/재시도", RGBColor(237, 248, 255))
    add_timeline_block(s, Inches(6.0), Inches(2.3), Inches(2.3), Inches(1.2), "2026-02-15", "Phase2 완료\n(checkpoint 확정)", RGBColor(237, 249, 240))
    add_timeline_block(s, Inches(8.6), Inches(2.3), Inches(2.3), Inches(1.2), "2026-02-15", "정량 비교\nTeacher 대비 53.9%", RGBColor(255, 247, 231))
    add_timeline_block(s, Inches(11.2), Inches(2.3), Inches(1.3), Inches(1.2), "2026-02-15", "Release\nTag", RGBColor(255, 235, 235))
    add_bullets(s, Inches(0.9), Inches(4.2), Inches(11.8), Inches(2.2), ["이번 사이클은 학습 완료 후 문서/릴리즈 정리까지 종료 처리"], size=20)

    # 6) 최종 상태
    s = prs.slides.add_slide(blank)
    add_title(s, "최종 상태 요약", "체크포인트/릴리즈/저장소 상태")
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(12.0),
        Inches(4.8),
        [
            "최종 distillation ckpt: checkpoints/distillation/phase2_epoch2_step24417.pt",
            "최종 merged ckpt: checkpoints/final/student_phase2_video_merged_20260215_161400.pt",
            "quantized ckpt: checkpoints/quantized/quantized_int8_int4.pt",
            "릴리즈 태그: v2026.02.15-project-close",
            "브랜치 정리 완료: main only",
        ],
        size=18,
    )

    # 7) 아키텍처 개요
    s = prs.slides.add_slide(blank)
    add_title(s, "모델 아키텍처 개요", "EfficientSAM3 Student 상위 블록")
    add_box(s, Inches(0.8), Inches(2.2), Inches(2.1), Inches(1.0), "입력 이미지", fill=RGBColor(250, 250, 255))
    add_box(s, Inches(0.8), Inches(3.7), Inches(2.1), Inches(1.0), "텍스트 프롬프트", fill=RGBColor(250, 250, 255))
    add_arrow(s, Inches(3.0), Inches(2.45), Inches(0.8), Inches(0.5))
    add_arrow(s, Inches(3.0), Inches(3.95), Inches(0.8), Inches(0.5))
    add_box(s, Inches(3.9), Inches(1.9), Inches(2.4), Inches(1.1), "Vision Encoder", fill=RGBColor(237, 248, 255))
    add_box(s, Inches(3.9), Inches(3.5), Inches(2.4), Inches(1.1), "Text Encoder\n(MobileCLIP)", fill=RGBColor(237, 248, 255))
    add_arrow(s, Inches(6.5), Inches(2.65), Inches(0.8), Inches(0.5))
    add_arrow(s, Inches(6.5), Inches(3.75), Inches(0.8), Inches(0.5))
    add_box(s, Inches(7.4), Inches(2.3), Inches(2.6), Inches(1.4), "DETR Decoder\n+ Query Matching", fill=RGBColor(237, 249, 240), border=SUCCESS)
    add_arrow(s, Inches(10.2), Inches(2.8), Inches(0.8), Inches(0.5))
    add_box(s, Inches(11.1), Inches(2.0), Inches(1.9), Inches(2.0), "출력\n- masks\n- boxes\n- iou\n- presence", fill=RGBColor(255, 247, 231), border=WARNING)

    # 8) 데이터 흐름
    s = prs.slides.add_slide(blank)
    add_title(s, "데이터 흐름", "입력부터 예측까지")
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(12.0),
        Inches(4.8),
        [
            "1) 이미지 전처리: student 입력 크기(504)로 정규화",
            "2) 텍스트 인코딩: student tokenizer(open_clip)로 프롬프트 토큰화",
            "3) 비전/텍스트 특징 결합 후 DETR 디코더에서 query 기반 예측 생성",
            "4) mask/box/logit/presence를 동시에 산출해 손실 항목별 최적화",
            "5) 평가에서는 mIoU, Prompt sensitivity, Inference latency를 함께 추적",
        ],
        size=20,
    )

    # 9) 지식 증류 원리
    s = prs.slides.add_slide(blank)
    add_title(s, "지식 증류 원리", "Teacher의 표현을 Student가 따라가도록 학습")
    add_box(s, Inches(1.0), Inches(2.1), Inches(4.7), Inches(1.4), "Teacher (SAM3)\n고품질 예측/특징 제공", fill=RGBColor(255, 238, 238), border=ERROR)
    add_arrow(s, Inches(5.9), Inches(2.45), Inches(1.1), Inches(0.5))
    add_box(s, Inches(7.2), Inches(2.1), Inches(4.7), Inches(1.4), "Student (EfficientSAM3)\n경량 구조로 Teacher를 근사", fill=RGBColor(237, 248, 255))
    add_bullets(
        s,
        Inches(1.0),
        Inches(4.1),
        Inches(11.2),
        Inches(2.2),
        [
            "핵심: 예측 결과뿐 아니라 중간 표현 정렬까지 학습",
            "효과: 모바일 친화 모델에서도 품질 하락을 완화",
            "제한: teacher 대비 품질 격차를 줄이려면 학습량/데이터 품질이 중요",
        ],
        size=19,
    )

    # 10) 학습 단계
    s = prs.slides.add_slide(blank)
    add_title(s, "학습 단계", "Phase별 역할 분리")
    add_kpi_card(s, Inches(0.8), Inches(2.0), Inches(3.8), Inches(1.6), "Phase 1", "Feature Alignment", "표현 정렬 및 안정화", color=RGBColor(237, 248, 255))
    add_kpi_card(s, Inches(4.9), Inches(2.0), Inches(3.8), Inches(1.6), "Phase 2", "Output Refinement", "mIoU 직접 개선 구간", color=RGBColor(237, 249, 240))
    add_kpi_card(s, Inches(9.0), Inches(2.0), Inches(3.4), Inches(1.6), "Post", "Video/Quant", "병합/양자화/배포 준비", color=RGBColor(255, 247, 231))
    add_bullets(s, Inches(0.9), Inches(4.2), Inches(11.8), Inches(2.2), ["이번 사이클은 Phase2 완료 후 품질 회복 및 문서화까지 종료"], size=20)

    # 11) 코드 변경
    s = prs.slides.add_slide(blank)
    add_title(s, "핵심 코드 변경사항", "문제 해결에 직접 기여한 수정")
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(12.0),
        Inches(4.8),
        [
            "train_distill.py: Phase1 체크포인트 자동 선택 로직을 step 기준으로 정렬하도록 보정",
            "scripts/visual_eval_student.py: teacher 로딩 시 resize_teacher_rope import 누락 수정",
            "README/docs: 목적/성과/종료 상태/배포 체크리스트 문서화",
            "릴리즈 태그/노트로 프로젝트 종료 지점을 명확히 고정",
        ],
        size=19,
    )

    # 12) 재학습 성과
    s = prs.slides.add_slide(blank)
    add_title(s, "재학습 정량 성과", "Before vs After (num_val=200)")
    add_bar_comparison(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(5.7),
        Inches(3.7),
        ["Before", "After"],
        [0.0570, 0.0898],
        [RGBColor(160, 160, 160), ACCENT],
        "mIoU",
    )
    add_bar_comparison(
        s,
        Inches(6.8),
        Inches(2.0),
        Inches(5.7),
        Inches(3.7),
        ["Before", "After"],
        [0.1479, 0.1343],
        [RGBColor(160, 160, 160), SUCCESS],
        "Prompt Sensitivity (낮을수록 좋음)",
    )
    add_bullets(s, Inches(0.9), Inches(5.9), Inches(11.8), Inches(1.0), ["mIoU +57.46%, Prompt sensitivity -9.17%"], size=18)

    # 13) Teacher 대비
    s = prs.slides.add_slide(blank)
    add_title(s, "Teacher 대비 위치", "Student 추종률 확인")
    add_bar_comparison(
        s,
        Inches(1.2),
        Inches(2.0),
        Inches(10.6),
        Inches(3.8),
        ["Student", "Teacher"],
        [0.0787, 0.1459],
        [ACCENT, ERROR],
        "mIoU Snapshot",
    )
    add_bullets(
        s,
        Inches(1.1),
        Inches(5.9),
        Inches(11.0),
        Inches(1.0),
        ["Student/Teacher = 0.5392 (약 53.9%) -> 추가 사이클 여지 큼"],
        size=18,
    )

    # 14) 양자화 결과
    s = prs.slides.add_slide(blank)
    add_title(s, "양자화 결과", "품질 보존성과 한계")
    add_bullets(
        s,
        Inches(0.9),
        Inches(2.0),
        Inches(11.8),
        Inches(4.8),
        [
            "FP16 mIoU: 0.0898",
            "int8_int4 mIoU: 0.0879 (drop 0.0019)",
            "의미: 품질 손실이 작아 양자화 방향은 유효",
            "주의: int4 단독 모드는 mslk 의존성 이슈로 실패",
        ],
        size=20,
    )

    # 15) 정성 결과 이미지
    s = prs.slides.add_slide(blank)
    add_title(s, "정성 시각 결과", "프롬프트별 마스크 출력 예시")
    imgs = [
        (ROOT / "docs" / "visual_qa_person.png", "person"),
        (ROOT / "docs" / "visual_qa_car.png", "car"),
        (ROOT / "docs" / "visual_qa_building.png", "building"),
        (ROOT / "docs" / "visual_qa_segment_everything.png", "segment everything"),
    ]
    positions = [
        (Inches(0.8), Inches(1.9)),
        (Inches(6.9), Inches(1.9)),
        (Inches(0.8), Inches(4.3)),
        (Inches(6.9), Inches(4.3)),
    ]
    for (path, cap), (lx, ty) in zip(imgs, positions):
        add_picture_or_placeholder(s, path, lx, ty, Inches(5.6), Inches(2.1), cap)

    # 16) 배포/운영 관점
    s = prs.slides.add_slide(blank)
    add_title(s, "배포/운영 관점 정리", "현재는 종료 정리 단계, 공개 배포는 추가 검증 필요")
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(12.0),
        Inches(4.8),
        [
            "완료: 코드/문서/태그/릴리즈까지 프로젝트 종료 절차 수행",
            "남은 과제: on-device KPI (size/latency/fps/NPU) 실측",
            "판단: 제한적 베타 가능, 대규모 공개 품질은 후속 사이클 권장",
            "정리: 이어학습 가능한 최소 자산(checkpoint + dataset) 보존",
        ],
        size=19,
    )

    # 17) 결론
    s = prs.slides.add_slide(blank)
    add_title(s, "결론 및 다음 사이클 제안", "이번 프로젝트는 '회복 + 정리'까지 완료")
    add_bullets(
        s,
        Inches(0.9),
        Inches(2.0),
        Inches(11.8),
        Inches(4.8),
        [
            "성과: 구조 안정화와 재학습으로 품질 회복 경로 확보",
            "한계: teacher 대비 품질 격차와 presence 지표 문제는 남아 있음",
            "다음: Phase2 연장 + negative sample + on-device KPI 실측",
            "상태: v2026.02.15-project-close 기준으로 아카이브 가능",
        ],
        size=21,
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    return OUT_PATH


def main():
    out = build_presentation()
    print(f"saved={out}")


if __name__ == "__main__":
    main()
